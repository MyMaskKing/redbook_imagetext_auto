import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import pandas as pd
from pptx import Presentation as PptxPresentation
import os
import time
from pptx.util import Pt
import traceback
from copy import deepcopy
from PIL import Image, ImageDraw, ImageFont, ImageTk
import io
from pptx.enum.shapes import MSO_SHAPE_TYPE
import win32com.client
import pythoncom
from spire.presentation import Presentation
import comtypes.client

class ModernButton(tk.Button):
    def __init__(self, master, **kwargs):
        # 提取自定义颜色参数
        self.start_color = kwargs.pop('start_color', '#FF4D6D') if 'start_color' in kwargs else '#FF4D6D'
        self.end_color = kwargs.pop('end_color', '#FF8FA3') if 'end_color' in kwargs else '#FF8FA3'
        
        # 设置基本配置
        kwargs.update({
            'background': self.start_color,
            'foreground': 'white',
            'font': ('Microsoft YaHei UI', 10),
            'borderwidth': 0,
            'activebackground': self.end_color,
            'activeforeground': 'white',
            'padx': 15,
            'pady': 8,
            'cursor': 'hand2',
            'relief': 'flat'
        })
        
        super().__init__(master, **kwargs)
        
        # 圆角效果（使用Canvas实现）
        self.canvas = tk.Canvas(self, width=20, height=20, bg=self.start_color, 
                              highlightthickness=0)
        self.canvas.create_arc(0, 0, 20, 20, start=90, extent=90, fill=self.start_color)
        
        self.bind('<Enter>', self.on_enter)
        self.bind('<Leave>', self.on_leave)

    def on_enter(self, e):
        self.config(background=self.end_color)
        self.canvas.config(bg=self.end_color)

    def on_leave(self, e):
        self.config(background=self.start_color)
        self.canvas.config(bg=self.start_color)

class PPTGeneratorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("小红书图文批量制作工具")
        
        # 设置窗口大小和位置
        window_width = 800
        window_height = 600
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        center_x = int(screen_width/2 - window_width/2)
        center_y = int(screen_height/2 - window_height/2)
        self.root.geometry(f'{window_width}x{window_height}+{center_x}+{center_y}')
        
        # 设置小红书风格主题（浅粉色背景）
        self.root.configure(bg='#FFF0F5')
        
        # 创建主容器
        main_container = tk.Frame(root, bg='#FFF0F5')
        main_container.pack(fill=tk.BOTH, expand=True)
        
        # 创建画布和滚动条
        self.canvas = tk.Canvas(main_container, bg='#FFF0F5', highlightthickness=0)
        self.scrollbar = ttk.Scrollbar(main_container, orient="vertical", command=self.canvas.yview)
        
        # 创建主框架
        self.main_frame = tk.Frame(self.canvas, bg='#FFF0F5')
        
        # 配置画布滚动区域
        self.canvas_frame = self.canvas.create_window((0, 0), window=self.main_frame, anchor="nw")
        
        # 绑定画布和滚动条
        self.canvas.configure(yscrollcommand=self.scrollbar.set)
        
        # 布局画布和滚动条
        self.canvas.pack(side="left", fill="both", expand=True, padx=(20, 0))  # 添加左边距
        self.scrollbar.pack(side="right", fill="y")
        
        # 绑定事件
        self.canvas.bind_all("<MouseWheel>", self._on_mousewheel)
        self.main_frame.bind("<Configure>", self._on_frame_configure)
        self.canvas.bind("<Configure>", self._on_canvas_configure)
        
        # 标题（使用小红书logo颜色）
        title_label = tk.Label(
            self.main_frame,
            text="小红书图文批量制作工具",
            font=('Microsoft YaHei UI', 24, 'bold'),
            fg='#FF2442',
            bg='#FFF0F5'
        )
        title_label.pack(pady=(0, 30))
        
        # 文件选择区域
        self.create_file_frame()
        
        # 设置区域
        self.create_settings_frame()
        
        # 添加尺寸设置区域
        self.create_scale_frame()
        
        # 进度条区域
        self.create_progress_frame()
        
        # WPS提示
        self.create_wps_notice()
        
        # 生成按钮
        self.create_generate_button()

    def _on_mousewheel(self, event):
        """处理鼠标滚轮事件"""
        self.canvas.yview_scroll(int(-1*(event.delta/120)), "units")

    def _on_frame_configure(self, event=None):
        """更新画布的滚动区域"""
        self.canvas.configure(scrollregion=self.canvas.bbox("all"))

    def _on_canvas_configure(self, event):
        """当画布大小改变时调整框架宽度"""
        # 设置框架宽度以匹配画布
        self.canvas.itemconfig(self.canvas_frame, width=event.width)

    def create_file_frame(self):
        file_frame = tk.LabelFrame(
            self.main_frame,
            text="文件选择",
            font=('Microsoft YaHei UI', 12, 'bold'),
            fg='#FF2442',
            bg='#FFFFFF',
            padx=20,
            pady=20,
            relief='flat',
        )
        file_frame.pack(fill=tk.X, pady=(0, 20), padx=20)

        # 配置列的权重
        file_frame.grid_columnconfigure(1, weight=1)  # 让输入框列自动扩展

        # 为每个选择按钮设置不同的渐变色
        self.ppt_path = self.create_file_entry(
            file_frame, "PPT模板:", self.select_ppt, 0,
            start_color='#FF4D6D', end_color='#FF8FA3'
        )
        
        self.excel_path = self.create_file_entry(
            file_frame, "Excel文件:", self.select_excel, 1,
            start_color='#FF6B6B', end_color='#FFA5A5'
        )
        
        self.save_path = self.create_file_entry(
            file_frame, "保存位置:", self.select_save_path, 2,
            start_color='#FF8882', end_color='#FFACAC'
        )

    def create_file_entry(self, parent, label_text, command, row, start_color, end_color):
        # 标签
        label = tk.Label(
            parent,
            text=label_text,
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            width=10,  # 固定标签宽度
            anchor='e'  # 右对齐
        )
        label.grid(row=row, column=0, sticky='e', pady=10, padx=(0, 10))
        
        # 输入框
        var = tk.StringVar()
        entry = tk.Entry(
            parent,
            textvariable=var,
            font=('Microsoft YaHei UI', 10),
            bg='#F8F8F8',
            fg='#333333',
            insertbackground='#666666',
            relief='flat',
            highlightthickness=1,
            highlightcolor='#FF4D6D',
            highlightbackground='#E0E0E0'
        )
        entry.grid(row=row, column=1, sticky='ew', padx=10)
        
        # 按钮
        button = ModernButton(
            parent,
            text="选择",
            command=command,
            start_color=start_color,
            end_color=end_color,
            width=8  # 固定按钮宽度
        )
        button.grid(row=row, column=2, padx=(0, 10))
        
        return var

    def create_settings_frame(self):
        settings_frame = tk.LabelFrame(
            self.main_frame,
            text="设置",
            font=('Microsoft YaHei UI', 12, 'bold'),
            fg='#FF2442',
            bg='#FFFFFF',
            padx=20,
            pady=20,
            relief='flat'
        )
        settings_frame.pack(fill=tk.X, pady=(0, 20), padx=20)

        # 使用Grid布局
        settings_frame.grid_columnconfigure(1, weight=1)
        
        # 创建两个容器框架，分别用于标题设置和标题处理
        title_container = tk.Frame(settings_frame, bg='#FFFFFF')
        title_container.grid(row=0, column=0, columnspan=2, sticky='ew', pady=(0, 10))
        
        process_container = tk.Frame(settings_frame, bg='#FFFFFF')
        process_container.grid(row=1, column=0, columnspan=2, sticky='ew')
        
        # 标题设置
        tk.Label(
            title_container,
            text="标题设置：",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            width=10,
            anchor='e'
        ).pack(side=tk.LEFT, padx=(0, 10))

        # 第一组单选按钮
        self.radio_var1 = tk.StringVar(value="option1")
        tk.Radiobutton(
            title_container,
            text="包含标题",
            variable=self.radio_var1,
            value="option1",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            activebackground='#FFE4E8',
            selectcolor='#FF4D6D'
        ).pack(side=tk.LEFT, padx=10)

        tk.Radiobutton(
            title_container,
            text="只有正文",
            variable=self.radio_var1,
            value="option2",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            activebackground='#FFE4E8',
            selectcolor='#FF4D6D'
        ).pack(side=tk.LEFT, padx=10)

        # 标题处理
        tk.Label(
            process_container,
            text="标题处理：",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            width=10,
            anchor='e'
        ).pack(side=tk.LEFT, padx=(0, 10))

        # 第二组单选按钮
        self.radio_var2 = tk.StringVar(value="option1")
        tk.Radiobutton(
            process_container,
            text="每页不同",
            variable=self.radio_var2,
            value="option1",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            activebackground='#FFE4E8',
            selectcolor='#FF4D6D'
        ).pack(side=tk.LEFT, padx=10)

        tk.Radiobutton(
            process_container,
            text="统一标题",
            variable=self.radio_var2,
            value="option2",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            activebackground='#FFE4E8',
            selectcolor='#FF4D6D'
        ).pack(side=tk.LEFT, padx=10)

    def create_scale_frame(self):
        scale_frame = tk.LabelFrame(
            self.main_frame,
            text="图片尺寸设置",
            font=('Microsoft YaHei UI', 12, 'bold'),
            fg='#FF2442',
            bg='#FFFFFF',
            padx=20,
            pady=20,
            relief='flat'
        )
        scale_frame.pack(fill=tk.X, pady=(0, 20), padx=20)

        # 使用容器来居中对齐内容
        container = tk.Frame(scale_frame, bg='#FFFFFF')
        container.pack(expand=True)
        
        # 宽度设置
        tk.Label(
            container,
            text="宽度:",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            width=6,
            anchor='e'
        ).grid(row=0, column=0, padx=(0, 5))
        
        self.width_var = tk.StringVar(value="1920")
        tk.Entry(
            container,
            textvariable=self.width_var,
            font=('Microsoft YaHei UI', 10),
            bg='#F8F8F8',
            fg='#333333',
            width=8,
            justify='center'
        ).grid(row=0, column=1, padx=5)
        
        # 高度设置
        tk.Label(
            container,
            text="高度:",
            font=('Microsoft YaHei UI', 10),
            fg='#333333',
            bg='#FFFFFF',
            width=6,
            anchor='e'
        ).grid(row=0, column=2, padx=(20, 5))
        
        self.height_var = tk.StringVar(value="1080")
        tk.Entry(
            container,
            textvariable=self.height_var,
            font=('Microsoft YaHei UI', 10),
            bg='#F8F8F8',
            fg='#333333',
            width=8,
            justify='center'
        ).grid(row=0, column=3, padx=5)

    def create_progress_frame(self):
        progress_frame = tk.Frame(self.main_frame, bg='#FFF0F5')
        progress_frame.pack(fill=tk.X, pady=(0, 20), padx=20)
        
        # 创建圆角进度条背景
        progress_bg = tk.Canvas(
            progress_frame,
            height=20,
            bg='#FFE4E8',
            highlightthickness=0
        )
        progress_bg.pack(fill=tk.X, padx=2)
        
        # 创建进度条
        self.progress_canvas = tk.Canvas(
            progress_bg,
            height=16,
            bg='#FFE4E8',
            highlightthickness=0
        )
        self.progress_canvas.place(relx=0.01, rely=0.5, relwidth=0.98, anchor='w')
        
        # 初始化进度变量
        self.progress_var = tk.DoubleVar(value=0)
        
        # 创建进度文本标签
        self.progress_label = tk.Label(
            progress_frame,
            text="准备就绪",
            font=('Microsoft YaHei UI', 10),
            fg='#666666',
            bg='#FFF0F5'
        )
        self.progress_label.pack(pady=10)

        # 绑定进度变量更新事件
        self.progress_var.trace_add('write', self._update_progress_bar)

    def _update_progress_bar(self, *args):
        """更新进度条显示"""
        progress = self.progress_var.get()
        width = self.progress_canvas.winfo_width()
        filled_width = int(width * (progress / 100))
        
        # 清除原有内容
        self.progress_canvas.delete('progress')
        
        # 绘制圆角进度条
        if filled_width > 0:
            # 计算圆角矩形的坐标
            x1, y1 = 0, 0
            x2, y2 = filled_width, 16
            radius = 8  # 圆角半径
            
            # 创建圆角矩形路径
            self.progress_canvas.create_polygon(
                x1+radius, y1,
                x2-radius, y1,
                x2, y1,
                x2, y2,
                x2-radius, y2,
                x1+radius, y2,
                x1, y2,
                x1, y1,
                fill='#FF4D6D',
                smooth=True,
                tags='progress'
            )

    def create_wps_notice(self):
        notice_frame = tk.Frame(self.main_frame, bg='#FFF0F5')
        notice_frame.pack(fill=tk.X, pady=(0, 20))
        
        tk.Label(
            notice_frame,
            text="注意：本机必须安装WPS软件，否则无法生成图片",
            font=('Microsoft YaHei UI', 10),
            fg='#FF2442',
            bg='#FFE4E8',
            padx=15,
            pady=10
        ).pack(fill=tk.X)

    def create_generate_button(self):
        button_frame = tk.Frame(self.main_frame, bg='#FFF0F5')
        button_frame.pack(pady=30)
        
        self.generate_button = ModernButton(
            button_frame,
            text="开始生成",
            command=self.generate_ppt,
            width=20,
            height=2,
            start_color='#FF2442',  # 小红书主色调
            end_color='#FF4D6D'
        )
        self.generate_button.pack()

    def select_ppt(self):
        filename = filedialog.askopenfilename(
            title="选择PPT模板文件",
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        if filename:
            self.ppt_path.set(filename)

    def select_excel(self):
        filename = filedialog.askopenfilename(
            title="选择Excel数据文件",
            filetypes=[("Excel files", "*.xlsx;*.xls")]
        )
        if filename:
            self.excel_path.set(filename)

    def select_save_path(self):
        dirname = filedialog.askdirectory(
            title="选择保存文件夹"
        )
        if dirname:
            # 生成默认文件名（当前时间）
            current_time = time.strftime("%Y%m%d_%H%M%S")
            default_filename = f"小红书图文_{current_time}.pptx"
            # 组合完整的保存路径
            save_path = os.path.join(dirname, default_filename)
            self.save_path.set(save_path)

    def print_shape_info(self, ppt):
        for slide in ppt.slides:
            print("\n=== 幻灯片信息 ===")
            for shape in slide.shapes:
                print(f"形状名称: {shape.name}")
                print(f"形状类型: {shape.shape_type}")
                if hasattr(shape, "text"):
                    print(f"文本内容: {shape.text}")
                if hasattr(shape, "placeholder_format"):
                    print(f"占位符类型: {shape.placeholder_format.type}")
                print("---")

    def convert_ppt_to_images(self, ppt_path, batch_size=50):
        try:
            # 获取文件名（不含扩展名）作为文件夹名
            base_name = os.path.splitext(os.path.basename(ppt_path))[0]
            
            # 创建与PPT同名的文件夹
            ppt_dir = os.path.dirname(ppt_path)
            images_dir = os.path.join(ppt_dir, base_name)
            if not os.path.exists(images_dir):
                os.makedirs(images_dir)

            # 初始化COM
            pythoncom.CoInitialize()
            
            try:
                # 启动 WPS 应用程序
                wps = comtypes.client.CreateObject("KWPP.Application")
                wps.Visible = True  # 设置为可见，避免一些COM错误
                
                # 打开PPT文件
                ppt = wps.Presentations
                presentation = ppt.Open(os.path.abspath(ppt_path))  # 使用绝对路径
                
                slide_count = presentation.Slides.Count
                for batch_start in range(0, slide_count, batch_size):
                    for i in range(batch_start, min(batch_start + batch_size, slide_count)):
                        slide = presentation.Slides.Item(i + 1)  # 使用Item方法
                        # 生成输出文件路径
                        output_path = os.path.join(images_dir, f"{base_name}_第{i+1}页.jpg")
                        
                        # 导出当前幻灯片为JPG格式
                        slide.Export(output_path, "JPG")
                        print(f"已将幻灯片 {i + 1} 保存为 {output_path}")
                        
                        # 使用Pillow来提高清晰度（DPI）
                        img = Image.open(output_path)
                        img.save(output_path, "JPEG", quality=95, dpi=(300, 300))
                    
                    print(f"已处理第{batch_start + 1}到{min(batch_start + batch_size, slide_count)}张幻灯片")
                
                return images_dir
                
            finally:
                # 清理资源
                try:
                    if 'presentation' in locals():
                        presentation.Close()
                    if 'wps' in locals():
                        wps.Quit()
                except:
                    pass
                pythoncom.CoUninitialize()
                
        except Exception as e:
            raise Exception(f"转换图片时出错: {str(e)}")

    def update_progress(self, value, message):
        """更新进度条和进度信息"""
        self.progress_var.set(value)
        self.progress_label.config(text=message)
        self.root.update()

    def generate_ppt(self):
        try:
            # 初始化进度
            self.update_progress(0, "开始处理...")
            
            # 验证文件路径
            if not self.ppt_path.get():
                messagebox.showerror("错误", "请选择PPT模板文件")
                return
            if not self.excel_path.get():
                messagebox.showerror("错误", "请选择Excel数据文件")
                return
            if not self.save_path.get():
                messagebox.showerror("错误", "请选择保存位置")
                return
            
            # 验证文件是否存在
            if not os.path.exists(self.ppt_path.get()):
                messagebox.showerror("错误", "PPT模板文件不存在")
                return
            if not os.path.exists(self.excel_path.get()):
                messagebox.showerror("错误", "Excel数据文件不存在")
                return
            
            # 验证保存路径的文件夹是否存在
            save_dir = os.path.dirname(self.save_path.get())
            if not os.path.exists(save_dir):
                messagebox.showerror("错误", "保存文件夹不存在")
                return

            self.update_progress(10, "读取Excel文件...")
            df = pd.read_excel(self.excel_path.get())
            
            self.update_progress(20, "加载PPT模板...")
            ppt = PptxPresentation(self.ppt_path.get())
            
            # 获取数据总行数用于计算进度
            total_rows = len(df.iloc[1:])
            
            # 获取单选按钮的值
            has_title = self.radio_var1.get() == "option1"  # 是否包含标题
            unified_title = self.radio_var2.get() == "option2"  # 是否统一标题
            
            # 获取模板第一页
            template_slide = ppt.slides[0]
            
            # 获取Excel的第一行（用于匹配PPT中的对象）
            headers = df.iloc[0]
            
            # 遍历Excel数据（从第二行开始）
            for index, row in df.iloc[1:].iterrows():
                progress = 20 + (index / total_rows * 40)  # 20-60%用于生成PPT
                self.update_progress(progress, f"正在处理第 {index} 行数据...")
                # 复制整个幻灯片
                new_slide = ppt.slides.add_slide(template_slide.slide_layout)
                
                # 删除新幻灯片中的默认形状
                for shape in new_slide.shapes:
                    element = shape._element
                    element.getparent().remove(element)
                
                # 从模板导入所有形状（包括格式和背景）
                for shape in template_slide.shapes:
                    el = shape.element
                    new_el = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                
                # 遍历所有形状并只更新文本内容
                for shape in new_slide.shapes:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    
                    # 获取shape的名称
                    shape_name = shape.name
                    print(f"当前形状名称: {shape_name}")  # 调试信息
                    
                    # 在Excel的列名中查找匹配的内容
                    for col, header_text in headers.items():
                        if str(col).strip() == shape_name.strip():
                            # 找到匹配的列，获取对应的内容
                            content = str(row[col])
                            print(f"匹配到列: {col}, 内容: {content}")  # 调试信息
                            
                            # 如果是标题且选择了"只有正文"，则跳过
                            if "标题" in shape_name and not has_title:
                                continue
                                
                            # 如果是标题且选择了"统一标题"
                            if "标题" in shape_name and unified_title:
                                content = str(df.iloc[1][col]) if index == 1 else shape.text_frame.text
                            
                            # 只更新文本内容，保持原有格式
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.runs:
                                    # 保持原有格式，只更新文本
                                    paragraph.runs[0].text = content
                                else:
                                    # 如果没有runs，创建新的并复制原有格式
                                    run = paragraph.add_run()
                                    run.text = content
                            break
            
            self.update_progress(60, "保存PPT文件...")
            ppt.save(self.save_path.get())
            
            self.update_progress(70, "转换为图片...")
            images_dir = self.convert_ppt_to_images(self.save_path.get())
            
            self.update_progress(90, "打开生成的文件...")
            try:
                os.startfile(self.save_path.get())
                os.startfile(images_dir)
            except Exception as open_error:
                print(f"打开文件失败: {str(open_error)}")
            
            self.update_progress(100, "处理完成！")
            messagebox.showinfo("成功", "PPT生成完成！图片已保存到images文件夹")
            
        except Exception as e:
            self.update_progress(0, "处理出错")
            messagebox.showerror("错误", f"生成过程中出现错误：{str(e)}\n{traceback.format_exc()}")

def main():
    root = tk.Tk()
    app = PPTGeneratorApp(root)
    root.mainloop()

if __name__ == "__main__":
    main() 