from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 打开PPT文件
prs = Presentation('./ppt/template.pptx')

# 遍历所有幻灯片
for slide in prs.slides:
    # 获取所有对象
    shapes = slide.shapes
    
    # 遍历所有对象
    for shape in shapes:
        print(shape.name)
        
        # 检查对象是否为文本框
        if hasattr(shape, 'text'):
            # 获取文本框的文本
            txBox = shape.text_frame
            shape.text_frame=1222
            if txBox:
                # 遍历所有段落
                for paragraph in txBox.paragraphs:
                    # 遍历所有运行
                    for run in paragraph.runs:
                        # 检查文本内容是否包含“标题”或“内容”
                        if '标题' in run.text:
                            # 向“标题”中写入文本
                            run.text = '新的标题文本'
                        elif '内容' in run.text:
                            # 向“内容”中写入文本
                            run.text = '新的内容文本'
                        print(run.text)

# 保存修改后的PPT文件
prs.save('./ppt/template.pptx')